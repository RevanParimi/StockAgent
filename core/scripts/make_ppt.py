"""
scripts/make_ppt.py
====================
Generates StockAgent presentation: outputs/StockAgent_Presentation.pptx
Run: python scripts/make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
from pathlib import Path

# ─── Color Palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_CARD      = RGBColor(0x12, 0x26, 0x3A)   # slightly lighter navy
CYAN         = RGBColor(0x00, 0xE5, 0xFF)   # bright cyan accent
GREEN        = RGBColor(0x00, 0xE6, 0x7A)   # emerald green
YELLOW       = RGBColor(0xFF, 0xD6, 0x00)   # gold / warning
RED          = RGBColor(0xFF, 0x4C, 0x4C)   # alert red
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GRAY         = RGBColor(0xAA, 0xBB, 0xCC)
CODE_BG      = RGBColor(0x08, 0x12, 0x1E)   # near-black for code blocks
ORANGE       = RGBColor(0xFF, 0x8C, 0x00)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


# ─── Helper Utilities ─────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def bg(slide, color=BG_DARK):
    add_rect(slide, 0, 0, W, H, color)


def text_box(slide, text, x, y, w, h, size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def heading(slide, text, y=Inches(0.25), color=CYAN, size=Pt(32)):
    text_box(slide, text, Inches(0.5), y, Inches(12.3), Inches(0.65),
             size=size, bold=True, color=color, align=PP_ALIGN.LEFT)


def subheading(slide, text, y=Inches(0.9), color=YELLOW, size=Pt(20)):
    text_box(slide, text, Inches(0.5), y, Inches(12.3), Inches(0.45),
             size=size, bold=True, color=color)


def divider(slide, y, color=CYAN):
    line = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def code_block(slide, text, x, y, w, h, font_size=Pt(11.5)):
    add_rect(slide, x, y, w, h, CODE_BG)
    txb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.1),
                                   w - Inches(0.24), h - Inches(0.2))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.name = "Consolas"
        run.font.color.rgb = GREEN


def badge(slide, label, x, y, bg_color=CYAN, text_color=BG_DARK):
    w, h = Inches(1.5), Inches(0.38)
    add_rect(slide, x, y, w, h, bg_color)
    text_box(slide, label, x, y + Inches(0.04), w, h - Inches(0.08),
             size=Pt(13), bold=True, color=text_color, align=PP_ALIGN.CENTER)


def bullet_block(slide, items, x, y, w, h, size=Pt(15), color=WHITE, gap=Inches(0.32)):
    """items = list of (text, color) or just strings"""
    cur_y = y
    for item in items:
        if isinstance(item, tuple):
            txt, col = item
        else:
            txt, col = item, color
        text_box(slide, txt, x, cur_y, w, gap, size=size, color=col)
        cur_y += gap


def two_col_table(slide, headers, rows, x, y, col_widths, row_h=Inches(0.38),
                  header_bg=CYAN, row_bg=BG_CARD, alt_bg=CODE_BG, font_size=Pt(13)):
    total_w = sum(col_widths)
    # header row
    cx = x
    for i, (h_text, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, cx, y, cw, row_h, header_bg)
        text_box(slide, h_text, cx + Inches(0.06), y + Inches(0.05),
                 cw - Inches(0.12), row_h - Inches(0.1),
                 size=font_size, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        cx += cw
    ry = y + row_h
    for ri, row in enumerate(rows):
        cx = x
        rbg = row_bg if ri % 2 == 0 else alt_bg
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, cx, ry, cw, row_h, rbg)
            text_box(slide, cell, cx + Inches(0.06), ry + Inches(0.04),
                     cw - Inches(0.12), row_h - Inches(0.08),
                     size=font_size, color=WHITE, align=PP_ALIGN.LEFT)
            cx += cw
        ry += row_h


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)

# gradient strip top
add_rect(slide, 0, 0, W, Inches(0.08), CYAN)

# big title
text_box(slide, "StockAgent",
         Inches(0.6), Inches(1.4), Inches(12), Inches(1.8),
         size=Pt(72), bold=True, color=CYAN, align=PP_ALIGN.LEFT)

text_box(slide, "Multi-Agent AI for Indian Stock Analysis",
         Inches(0.6), Inches(3.1), Inches(12), Inches(0.7),
         size=Pt(28), bold=False, color=WHITE, align=PP_ALIGN.LEFT)

text_box(slide, "28 Specialist Agents  ·  4 Market Sectors  ·  4 Languages  ·  Real-Time Data",
         Inches(0.6), Inches(3.85), Inches(12), Inches(0.5),
         size=Pt(18), color=YELLOW, align=PP_ALIGN.LEFT)

divider(slide, Inches(4.5), CYAN)

# pill badges
for i, (label, color) in enumerate([
    ("Python", GREEN), ("C++", ORANGE), ("TypeScript", CYAN), ("C#", RGBColor(0x68, 0x17, 0xA8))
]):
    badge(slide, label, Inches(0.6 + i * 1.65), Inches(4.75), color, WHITE if color != CYAN else BG_DARK)

text_box(slide, "LangGraph  ·  OpenRouter/Qwen3  ·  ChromaDB RAG  ·  FastAPI  ·  Quartz.NET",
         Inches(0.6), Inches(5.4), Inches(12), Inches(0.45),
         size=Pt(15), color=GRAY)

text_box(slide, "Revan Parimi  |  April 2026",
         Inches(0.6), Inches(6.6), Inches(5), Inches(0.4),
         size=Pt(14), color=GRAY)

add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), CYAN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Why Multi-Agent?  The Core Problem")
divider(slide, Inches(0.95))

# Left box — naive approach
add_rect(slide, Inches(0.5), Inches(1.1), Inches(5.5), Inches(5.8), BG_CARD,
         RED, Pt(2))
text_box(slide, "❌  Single LLM Prompt", Inches(0.7), Inches(1.25), Inches(5), Inches(0.5),
         size=Pt(18), bold=True, color=RED)
code_block(slide,
           'User: "Should I buy MARUTI?"\n\n'
           'LLM: "MARUTI is a leading auto\n'
           'company. Strong fundamentals.\n'
           'Consider EV transition risks.\n'
           'Diversify your portfolio."\n\n'
           '→ Shallow  →  Generic\n'
           '→  No data  →  Not actionable',
           Inches(0.6), Inches(1.8), Inches(5.3), Inches(3.5), font_size=Pt(12))
text_box(slide, "One prompt cannot go deep on\nTechnicals AND Fundamentals\nAND Policy AND Macro simultaneously",
         Inches(0.7), Inches(5.35), Inches(5.2), Inches(1.3),
         size=Pt(13), color=GRAY)

# Right box — multi-agent
add_rect(slide, Inches(6.4), Inches(1.1), Inches(6.4), Inches(5.8), BG_CARD,
         GREEN, Pt(2))
text_box(slide, "✅  StockAgent: 8 Specialists", Inches(6.6), Inches(1.25), Inches(6), Inches(0.5),
         size=Pt(18), bold=True, color=GREEN)

agents_list = [
    "📊  Fundamentals Agent    → P&L, margins, EPS trend",
    "📈  Pattern Analysis      → RSI, MACD, support/resistance",
    "🏭  Raw Materials         → steel, crude, aluminium",
    "🚗  Sales & Demand        → FADA dispatch, EV growth",
    "📰  Policy & Regulatory   → FAME subsidies, BS7/CAFE",
    "🏆  Competitive Intel     → EV market share, new models",
    "🌍  Risk & Macro          → INR/USD, repo rate, geopolitics",
    "💬  Sentiment             → news NLP, management tone",
]
for i, a in enumerate(agents_list):
    text_box(slide, a, Inches(6.6), Inches(1.85 + i * 0.54), Inches(6.0), Inches(0.5),
             size=Pt(13), color=WHITE)

text_box(slide, "Each agent forced to go DEEP\non its own domain → Signal Aggregator\nfuses all 8 views into one verdict",
         Inches(6.6), Inches(6.2), Inches(5.9), Inches(1.0),
         size=Pt(13), color=YELLOW, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — KEY CONCEPT: LLM AGENT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Key Concept 1 — What Is an LLM Agent?")
divider(slide, Inches(0.95))

# Pipeline boxes
boxes = [
    ("CONTEXT\n(Live Data)", BG_CARD, CYAN),
    ("PROMPT\n(Instructions)", BG_CARD, YELLOW),
    ("LLM CALL\n(Qwen3-235B)", BG_CARD, ORANGE),
    ("STRUCTURED\nJSON OUTPUT", BG_CARD, GREEN),
]
for i, (label, bg_c, ac) in enumerate(boxes):
    bx = Inches(0.5 + i * 3.1)
    add_rect(slide, bx, Inches(1.1), Inches(2.7), Inches(1.0), bg_c, ac, Pt(2))
    text_box(slide, label, bx + Inches(0.1), Inches(1.2), Inches(2.5), Inches(0.8),
             size=Pt(16), bold=True, color=ac, align=PP_ALIGN.CENTER)
    if i < 3:
        text_box(slide, "→", Inches(3.1 + i * 3.1), Inches(1.35), Inches(0.5), Inches(0.5),
                 size=Pt(28), bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Context detail
text_box(slide, "Context Sources:", Inches(0.5), Inches(2.3), Inches(3), Inches(0.4),
         size=Pt(14), bold=True, color=CYAN)
for i, src in enumerate(["yfinance: OHLCV, financials", "Serper: Google news snippets",
                          "Tavily: full policy documents", "ChromaDB RAG (optional)"]):
    text_box(slide, f"  • {src}", Inches(0.5), Inches(2.7 + i * 0.35), Inches(3), Inches(0.35),
             size=Pt(13), color=WHITE)

# JSON output example
code_block(slide,
           '{\n'
           '  "overall_score": 0.68,\n'
           '  "sub_scores": {\n'
           '    "price_cycle_position": 0.62,\n'
           '    "seasonal_pattern":     0.72,\n'
           '    "rsi_macd_bb":          0.65,\n'
           '    "breakout_support_zone":0.55,\n'
           '    "peer_correlation":     0.60\n'
           '  },\n'
           '  "key_positives": ["MACD bullish crossover",\n'
           '                    "Q3 festive tailwind"],\n'
           '  "key_risks": ["Near resistance at 12,800"],\n'
           '  "summary": "Mid-cycle, momentum building"\n'
           '}',
           Inches(3.7), Inches(2.2), Inches(5.5), Inches(4.6), font_size=Pt(12))

text_box(slide, "response_format={\"type\": \"json_object\"}\nguarantees parseable output.\nScores clamped to [0.0, 1.0].\nOn parse failure → score=0.5 (neutral fallback).",
         Inches(9.5), Inches(2.2), Inches(3.5), Inches(2.0),
         size=Pt(13), color=GRAY)

text_box(slide, "Retry: 3 attempts with exponential back-off\nModel: qwen/qwen3-235b-a22b via OpenRouter",
         Inches(9.5), Inches(4.5), Inches(3.5), Inches(1.0),
         size=Pt(13), color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — LANGGRAPH
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Key Concept 2 — LangGraph: The Orchestration Brain")
divider(slide, Inches(0.95))

# Left: concept explanation
text_box(slide, "LangGraph is a graph-based orchestration\nframework for LLM workflows.",
         Inches(0.5), Inches(1.05), Inches(5.5), Inches(0.8), size=Pt(16), color=WHITE)

concepts = [
    ("StateGraph", "Directed graph where nodes are functions\nand edges define flow.", CYAN),
    ("Send()", "Fan-out primitive — spawns N parallel\nbranches, each with its own state.", YELLOW),
    ("Reducer", "Fan-in merge — combines all parallel\nbranch results into one state dict.", ORANGE),
    ("RetryPolicy", "Automatic retry per node if agent\nfails (max_attempts=2 here).", GREEN),
    ("input_rail /\noutput_rail", "Guard nodes: validate input ticker,\nclamp scores 0-1, inject defaults.", RGBColor(0xFF, 0x6E, 0xFF)),
]
for i, (term, desc, color) in enumerate(concepts):
    by = Inches(1.85 + i * 1.02)
    add_rect(slide, Inches(0.5), by, Inches(5.5), Inches(0.95), BG_CARD, color, Pt(1))
    text_box(slide, term, Inches(0.65), by + Inches(0.05), Inches(2.0), Inches(0.42),
             size=Pt(15), bold=True, color=color)
    text_box(slide, desc, Inches(0.65), by + Inches(0.46), Inches(5.1), Inches(0.44),
             size=Pt(12.5), color=GRAY)

# Right: graph topology
code_block(slide,
           'graph.invoke({"ticker": "TCS"})\n'
           '      │\n'
           '  [resolve_ticker]  ← LLM temp=0 → StockQuery\n'
           '      │\n'
           '  [input_rail]  ← yfinance ticker check\n'
           '      │\n'
           '      ├── Send("run_agent", {agent: fund})   ─┐\n'
           '      ├── Send("run_agent", {agent: macro})   │ PARALLEL\n'
           '      ├── Send("run_agent", {agent: risk})    │ fan-out\n'
           '      ├── Send("run_agent", {agent: pattern}) │\n'
           '      └── Send("run_agent", {agent: ...N})  ──┘\n'
           '               │\n'
           '               │  _merge_dicts reducer collects all\n'
           '               │  AgentOutput objects into state\n'
           '               │\n'
           '  [aggregate]  ← conflict_rail if spread > 0.35\n'
           '      │           → LLM re-resolution\n'
           '      │\n'
           '     END   →  state["final_report"] = FinalReport',
           Inches(6.3), Inches(1.05), Inches(6.5), Inches(6.2), font_size=Pt(11))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — CONTEXTBUILDER
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Key Concept 3 — ContextBuilder: Data Routing Engine")
divider(slide, Inches(0.95))

# Routing diagram
code_block(slide,
           'ContextBuilder.build(agent_name, query, sector="bfsi")\n'
           '      │\n'
           '      ▼ Step 1: Sector-specific?\n'
           '      ├── _build_bfsi_fundamentals()  ← NPA/NIM/CASA\n'
           '      ├── _build_bfsi_risk()           ← slippage, concentration\n'
           '      ├── _build_macro_policy()        ← RBI repo, MPC stance\n'
           '      │\n'
           '      ▼ Step 2: Generic automobile fallback?\n'
           '      ├── _build_pattern_analysis()   ← yfinance OHLCV (any sector)\n'
           '      ├── _build_fundamentals()       ← quarterly P&L\n'
           '      │\n'
           '      ▼ Step 3: Last resort\n'
           '      └── _build_generic()  → "Stock: MARUTI | Date: 2026-04-18"',
           Inches(0.5), Inches(1.05), Inches(6.7), Inches(4.2), font_size=Pt(11))

text_box(slide, "Why 3 levels?",
         Inches(0.5), Inches(5.35), Inches(3), Inches(0.4),
         size=Pt(15), bold=True, color=YELLOW)
text_box(slide,
         "• pattern_analysis is pure yfinance → works for\n  ALL sectors via fallthrough\n"
         "• sector-specific builders write NPA queries,\n  not crude-oil queries\n"
         "• stub ensures pipeline NEVER crashes",
         Inches(0.5), Inches(5.75), Inches(6.5), Inches(1.5),
         size=Pt(13), color=WHITE)

# Right: data fetcher cost table
text_box(slide, "Data Fetchers", Inches(7.6), Inches(1.05), Inches(5), Inches(0.4),
         size=Pt(18), bold=True, color=CYAN)

two_col_table(slide,
    ["Fetcher", "Returns", "Cost"],
    [
        ["get_fundamentals_context()", "Revenue, EBITDA, margins (yfinance)", "Free"],
        ["get_technical_context()", "RSI, MACD, Bollinger Bands", "Free"],
        ["get_macro_context()", "INR/USD, crude, repo rate", "Free"],
        ["fetch_news_context()", "Google snippets via Serper", "2,500/mo free"],
        ["fetch_tavily_context()", "Full PDF/page extraction", "1,000/mo free"],
    ],
    Inches(7.6), Inches(1.5),
    [Inches(2.5), Inches(2.8), Inches(1.5)],
    row_h=Inches(0.46), font_size=Pt(11.5))

text_box(slide, "Why not give every agent ALL data?",
         Inches(7.6), Inches(4.5), Inches(5.5), Inches(0.4),
         size=Pt(15), bold=True, color=YELLOW)
text_box(slide,
         "Token limit = 2,048 tokens per agent response.\n"
         "Stuffing irrelevant data wastes tokens and\n"
         "dilutes the signal. Each agent gets only\n"
         "what IT needs for its specific analysis.",
         Inches(7.6), Inches(4.95), Inches(5.5), Inches(1.8),
         size=Pt(13), color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RAG PIPELINE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Key Concept 4 — RAG: Retrieval-Augmented Generation")
divider(slide, Inches(0.95))

text_box(slide,
         "RAG lets agents reason over YOUR documents (annual reports, SEBI filings,\n"
         "MNRE circulars) — not just LLM training data or live news.",
         Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.7),
         size=Pt(15), color=GRAY)

# Ingestion flow
text_box(slide, "INGESTION  (one-time, offline)", Inches(0.5), Inches(1.8),
         Inches(5.8), Inches(0.4), size=Pt(15), bold=True, color=GREEN)
code_block(slide,
           'python scripts/ingest_documents.py --dir data/\n\n'
           'PDF / TXT / MD\n'
           '  └─► _chunk_text(512 tokens, 64 overlap)\n'
           '        └─► sentence-transformers.embed_batch()\n'
           '              └─► ChromaDB.upsert()\n'
           '                   metadata: {ticker, doc_type,\n'
           '                              chunk_index, source}\n'
           '                   ID: MD5(chunk) → re-ingest = idempotent',
           Inches(0.5), Inches(2.25), Inches(5.8), Inches(3.4), font_size=Pt(11.5))

# Retrieval flow
text_box(slide, "RETRIEVAL  (per agent, runtime)", Inches(0.5), Inches(5.75),
         Inches(5.8), Inches(0.4), size=Pt(15), bold=True, color=CYAN)
code_block(slide,
           '_rag_retrieve(query="MARUTI NIM trend")\n'
           '  └─► embed(query) → single vector\n'
           '        └─► ChromaDB cosine similarity\n'
           '              filter: {ticker: "MARUTI"}\n'
           '              → top-5 chunks (score ≥ 0.75)\n'
           '        └─► [optional] CrossEncoder re-rank\n'
           '              → ordered by relevance\n'
           '  returns: "---".join(chunks) → agent prompt',
           Inches(0.5), Inches(6.2), Inches(5.8), Inches(1.05), font_size=Pt(11.5))

# Right: priority chain
text_box(slide, "Context Priority Chain", Inches(7.0), Inches(1.8),
         Inches(5.8), Inches(0.4), size=Pt(15), bold=True, color=CYAN)
for i, (label, color, desc) in enumerate([
    ("1.  RAG (ChromaDB)", GREEN, "Your own docs — highest precision"),
    ("2.  ContextBuilder (live)", CYAN, "yfinance + Serper + Tavily"),
    ("3.  Minimal stub", YELLOW, '"Stock: MARUTI | Date: ..." — never crashes'),
]):
    by = Inches(2.3 + i * 1.3)
    add_rect(slide, Inches(7.0), by, Inches(6.0), Inches(1.1), BG_CARD, color, Pt(2))
    text_box(slide, label, Inches(7.15), by + Inches(0.08), Inches(5.7), Inches(0.45),
             size=Pt(16), bold=True, color=color)
    text_box(slide, desc, Inches(7.15), by + Inches(0.55), Inches(5.7), Inches(0.45),
             size=Pt(13), color=GRAY)

text_box(slide,
         "ChromaDB stores embeddings locally — no cloud,\n"
         "no latency, full control. Cosine similarity finds\n"
         "semantically similar chunks even with different wording.",
         Inches(7.0), Inches(6.25), Inches(6.0), Inches(1.0),
         size=Pt(13), color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SIGNAL AGGREGATION + CONFLICT RESOLUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Key Concept 5 — Signal Aggregation & Conflict Resolution")
divider(slide, Inches(0.95))

# Left: weighted fusion example
text_box(slide, "Step 1: Weighted Fusion", Inches(0.5), Inches(1.1),
         Inches(6.0), Inches(0.4), size=Pt(16), bold=True, color=CYAN)
code_block(slide,
           'Agent              Score   ×   Weight   =   Weighted\n'
           '───────────────────────────────────────────────────\n'
           'sales_demand       0.72   ×   0.18     =   0.130\n'
           'raw_materials      0.60   ×   0.10     =   0.060\n'
           'fundamentals       0.68   ×   0.20     =   0.136\n'
           'pattern_analysis   0.63   ×   0.13     =   0.082\n'
           'policy_regulatory  0.65   ×   0.10     =   0.065\n'
           'competitive_intel  0.62   ×   0.10     =   0.062\n'
           'risk_macro         0.58   ×   0.15     =   0.087\n'
           'sentiment          0.70   ×   0.04     =   0.028\n'
           '───────────────────────────────────────────────────\n'
           'COMPOSITE                               =   0.650  → BUY',
           Inches(0.5), Inches(1.55), Inches(6.2), Inches(3.6), font_size=Pt(11.5))

# Verdict table
text_box(slide, "Score → Verdict Mapping", Inches(0.5), Inches(5.2),
         Inches(6.0), Inches(0.4), size=Pt(15), bold=True, color=YELLOW)
verdicts = [
    ("0.75 – 1.00", "🚀 STRONG BUY", GREEN),
    ("0.55 – 0.75", "✅ BUY", CYAN),
    ("0.40 – 0.55", "⚖️  NEUTRAL", YELLOW),
    ("0.20 – 0.40", "⚠️  SELL", ORANGE),
    ("0.00 – 0.20", "🔴 STRONG SELL", RED),
]
for i, (rng, verdict, col) in enumerate(verdicts):
    bx, by = Inches(0.5 + i * 1.24), Inches(5.65)
    add_rect(slide, bx, by, Inches(1.18), Inches(1.3), BG_CARD, col, Pt(2))
    text_box(slide, verdict, bx + Inches(0.05), by + Inches(0.08),
             Inches(1.1), Inches(0.55), size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
    text_box(slide, rng, bx + Inches(0.05), by + Inches(0.7),
             Inches(1.1), Inches(0.45), size=Pt(12), color=GRAY, align=PP_ALIGN.CENTER)

# Right: conflict resolution
text_box(slide, "Step 2: Conflict Detection", Inches(7.0), Inches(1.1),
         Inches(6.0), Inches(0.4), size=Pt(16), bold=True, color=RED)
code_block(slide,
           '# Flag conflict if agents disagree by ≥0.30\n'
           'if abs(score_A - score_B) >= 0.30:\n'
           '    conflict_flags.append((A, B))\n\n'
           'Example conflict:\n'
           '  fundamentals  = 0.75  (BUY range)\n'
           '  risk_macro    = 0.35  (SELL range)\n'
           '  |0.75 - 0.35| = 0.40 ≥ 0.30  → CONFLICT',
           Inches(7.0), Inches(1.55), Inches(6.0), Inches(2.6), font_size=Pt(11.5))

text_box(slide, "Step 3: LLM Resolution (not arithmetic!)", Inches(7.0), Inches(4.3),
         Inches(6.0), Inches(0.4), size=Pt(15), bold=True, color=YELLOW)
text_box(slide,
         "Signal Aggregator sends all 8 scores + weights +\n"
         "conflict list to Qwen3. The LLM reasons:\n\n"
         "  \"Fundamentals strong (₹18K Cr revenue growth)\n"
         "   but macro bearish (RBI tightening, crude +12%).\n"
         "   Macro risk is near-term; fundamentals are structural.\n"
         "   Adjusted score: 0.62 → BUY with caution.\"\n\n"
         "A formula (discard outliers) loses information.\n"
         "Reasoning determines WHICH signal is more reliable now.",
         Inches(7.0), Inches(4.75), Inches(6.0), Inches(2.5),
         size=Pt(13), color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — MACRO CACHE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Key Concept 6 — Macro Cache: Smart API Budget Optimization")
divider(slide, Inches(0.95))

text_box(slide,
         "Serper free tier = 2,500 queries/month. Analysing 5 tickers/day × 3 sectors × 22 days = 330 sector-level calls wasted\n"
         "asking the SAME Fed rate / RBI repo question for every ticker. Macro Cache solves this.",
         Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.7), size=Pt(14), color=GRAY)

# Architecture diagram
code_block(slide,
           'main.py: _micro_search_loop()  [background daemon, every 4 hours]\n'
           '    │\n'
           '    ├── fetch_news_context(automobile_queries[:2])\n'
           '    │       → set_macro_cache("automobile", text)  ← key = "automobile"\n'
           '    │\n'
           '    ├── fetch_news_context(bfsi_queries[:2])\n'
           '    │       → set_macro_cache("bfsi", text)        ← key = "bfsi"\n'
           '    │\n'
           '    └── fetch_news_context(it_queries[:2])\n'
           '            → set_macro_cache("it", text)          ← key = "it"\n\n'
           'Per-analysis (e.g. analysing HDFCBANK):\n'
           '    ContextBuilder._build_macro_policy()\n'
           '        cached = get_macro_cache("bfsi")   ← TTL = 2 hours\n'
           '        if cached: return cached           ← SAVES 3 Serper calls\n'
           '        else: fetch Serper (cold start only)',
           Inches(0.5), Inches(1.8), Inches(7.5), Inches(4.4), font_size=Pt(11.5))

# Budget table
text_box(slide, "Budget Math", Inches(8.3), Inches(1.8),
         Inches(4.7), Inches(0.4), size=Pt(16), bold=True, color=YELLOW)
two_col_table(slide,
    ["Item", "Calls/month"],
    [
        ["Micro loop (3 sectors × 2 q × 6 cycles × 30d)", "1,080"],
        ["Saved (3 calls × 5 tickers × 3 sectors × 22d)", "−990"],
        ["Net loop overhead", "+90"],
    ],
    Inches(8.3), Inches(2.25),
    [Inches(3.5), Inches(1.2)],
    row_h=Inches(0.46), font_size=Pt(12))

text_box(slide, "Queries cached per sector:", Inches(8.3), Inches(4.0),
         Inches(4.7), Inches(0.4), size=Pt(14), bold=True, color=CYAN)
queries = [
    ("automobile", "Nifty Auto outlook, crude/steel/EV policy"),
    ("bfsi", "RBI MPC repo, credit growth, NPA quality"),
    ("it", "US IT spend, Fed rate, H1B, GenAI demand"),
]
for i, (sec, q) in enumerate(queries):
    by = Inches(4.45 + i * 0.78)
    add_rect(slide, Inches(8.3), by, Inches(4.7), Inches(0.72), BG_CARD, CYAN, Pt(1))
    text_box(slide, sec.upper(), Inches(8.45), by + Inches(0.04), Inches(1.0), Inches(0.32),
             size=Pt(13), bold=True, color=CYAN)
    text_box(slide, q, Inches(8.45), by + Inches(0.38), Inches(4.4), Inches(0.3),
             size=Pt(12), color=GRAY)

text_box(slide,
         "RE sector EXCLUDED — MNRE auctions, DISCOM\n"
         "payments, curtailment vary per company,\n"
         "not sector-wide. No cache benefit.",
         Inches(0.5), Inches(6.35), Inches(7.5), Inches(1.0),
         size=Pt(13), color=ORANGE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MULTI-LANGUAGE ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "System Architecture — 4 Languages, 1 Pipeline")
divider(slide, Inches(0.95))

lang_blocks = [
    ("Python  :8000", BG_CARD, GREEN,
     "8+ LLM agents  ·  LangGraph graphs  ·  FastAPI\n"
     "ContextBuilder  ·  RAG pipeline  ·  APScheduler\n"
     "POST /analyse  ·  WS /ws/stream  ·  GET /history"),
    ("C++  (in-process)", BG_CARD, ORANGE,
     "RSI(14), MACD(12,26,9), Bollinger Bands(20,2σ)\n"
     "pybind11 .pyd extension: stockindicators\n"
     "Pure-Python fallback if .pyd absent"),
    ("TypeScript  :3000/:3001", BG_CARD, CYAN,
     "Express REST proxy → :8000 (POST /api/analyse)\n"
     "WebSocket hub → :8000/ws/stream rebroadcast\n"
     "GET /api/schedule → :5000/scheduler/status"),
    ("C#  :5000", BG_CARD, RGBColor(0xA0, 0x60, 0xFF),
     "Quartz.NET cron: 8:30am IST weekdays\n"
     "EF Core → SQL Server (ScoreRecords table)\n"
     "POST :8000/analyse → save score → fire alert"),
]
for i, (title, bgc, ac, desc) in enumerate(lang_blocks):
    bx = Inches(0.35 + (i % 2) * 6.45)
    by = Inches(1.15 + (i // 2) * 2.85)
    add_rect(slide, bx, by, Inches(6.15), Inches(2.65), bgc, ac, Pt(2))
    text_box(slide, title, bx + Inches(0.15), by + Inches(0.12), Inches(5.8), Inches(0.5),
             size=Pt(20), bold=True, color=ac)
    text_box(slide, desc, bx + Inches(0.15), by + Inches(0.65), Inches(5.8), Inches(1.85),
             size=Pt(13.5), color=WHITE)

text_box(slide,
         "All language boundaries use JSON over HTTP — no shared memory, no cross-language FFI  "
         "(C++ is the only exception: in-process via pybind11 for zero-copy array passing).",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
         size=Pt(13.5), color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — 4 SECTOR GRAPHS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "4 Independent LangGraph Sector Graphs")
divider(slide, Inches(0.95))

sectors = [
    ("🚗  Automobile", "8 agents", "fundamentals 0.20", "FADA/Vahan dispatch\nCompetitive Intel (EV)\nFAME policy (Tavily)", GREEN),
    ("🏦  Banking / BFSI", "6 agents", "fundamentals 0.25", "NPA / NIM / CASA\nRBI MPC rate cycle\nBFSI macro cache", CYAN),
    ("💻  IT Sector", "8 agents", "fundamentals 0.25", "US tech spend + visa risk\nTranscript NLP (Tavily)\nIT macro cache", YELLOW),
    ("🌱  Renewable Energy", "6 agents", "fundamentals 0.30", "CUF / DSCR / EV/MW\nMNRE auctions (Tavily)\nRisk monitor (weight=0)", ORANGE),
]
for i, (title, n_agents, top_weight, signals, color) in enumerate(sectors):
    bx = Inches(0.3 + i * 3.22)
    add_rect(slide, bx, Inches(1.1), Inches(3.05), Inches(6.0), BG_CARD, color, Pt(2))
    text_box(slide, title, bx + Inches(0.1), Inches(1.2), Inches(2.85), Inches(0.55),
             size=Pt(15), bold=True, color=color)
    text_box(slide, n_agents, bx + Inches(0.1), Inches(1.78), Inches(2.85), Inches(0.38),
             size=Pt(13), bold=True, color=WHITE)
    text_box(slide, f"Top weight:\n{top_weight}", bx + Inches(0.1), Inches(2.18),
             Inches(2.85), Inches(0.7), size=Pt(12), color=GRAY)
    text_box(slide, "Unique signals:", bx + Inches(0.1), Inches(2.95),
             Inches(2.85), Inches(0.35), size=Pt(12), bold=True, color=color)
    text_box(slide, signals, bx + Inches(0.1), Inches(3.32),
             Inches(2.85), Inches(1.5), size=Pt(12), color=WHITE)

# Shared row
add_rect(slide, Inches(0.3), Inches(5.15), Inches(12.73), Inches(0.9), BG_CARD, GRAY, Pt(1))
text_box(slide, "All 4 graphs share:",
         Inches(0.5), Inches(5.2), Inches(2.0), Inches(0.4), size=Pt(14), bold=True, color=GRAY)
text_box(slide,
         "BaseAgent  ·  ContextBuilder (sector routing)  ·  macro_cache  ·  LangGraph node factories  ·  "
         "input_rail / output_rail / conflict_rail  ·  FinalReport schema",
         Inches(2.5), Inches(5.25), Inches(10.3), Inches(0.75), size=Pt(13), color=WHITE)

# Comparison table at bottom
two_col_table(slide,
    ["", "Automobile", "BFSI", "IT", "RE"],
    [
        ["Serper/run (warm)", "9", "12", "17", "15"],
        ["Tavily/run", "2", "0", "2", "2"],
        ["Macro cache", "✓ automobile", "✓ bfsi", "✓ it", "✗ (per-company)"],
    ],
    Inches(0.3), Inches(6.1),
    [Inches(2.0), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)],
    row_h=Inches(0.38), font_size=Pt(11.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DATA FLOW WALKTHROUGH
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "End-to-End: \"MARUTI\" → Investment Verdict")
divider(slide, Inches(0.95))

steps = [
    ('1', 'Ticker Resolution',
     'LLM call (temp=0.0)\n"Maruti Suzuki" → MARUTI.NS\nPydantic StockQuery object', CYAN),
    ('2', 'Context Assembly',
     'RAG → ContextBuilder → stub\nPer-agent: yfinance + Serper\n~440–2825 context tokens', YELLOW),
    ('3', 'Parallel Agents',
     'asyncio.gather × 8\nor ThreadPoolExecutor × 8\nAll agents fire simultaneously', ORANGE),
    ('4', 'LLM per Agent',
     'Qwen3-235B via OpenRouter\nJSON output: 5 sub-scores\nRetry 3× on failure', GREEN),
    ('5', 'Aggregation',
     'Weighted composite score\nConflict detection ≥0.30\nLLM resolves conflicts', CYAN),
    ('6', 'Final Verdict',
     'Score → STRONG BUY/BUY/\nNEUTRAL/SELL/STRONG SELL\nSaved to SQLite / SQL Server', YELLOW),
]
for i, (num, title, desc, color) in enumerate(steps):
    bx = Inches(0.35 + (i % 3) * 4.3)
    by = Inches(1.1 + (i // 3) * 2.85)
    add_rect(slide, bx, by, Inches(4.0), Inches(2.6), BG_CARD, color, Pt(2))
    # step number circle
    add_rect(slide, bx + Inches(0.1), by + Inches(0.1), Inches(0.55), Inches(0.55), color)
    text_box(slide, num, bx + Inches(0.1), by + Inches(0.1), Inches(0.55), Inches(0.55),
             size=Pt(18), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    text_box(slide, title, bx + Inches(0.75), by + Inches(0.12), Inches(3.1), Inches(0.5),
             size=Pt(15), bold=True, color=color)
    text_box(slide, desc, bx + Inches(0.15), by + Inches(0.72), Inches(3.7), Inches(1.7),
             size=Pt(13), color=WHITE)

text_box(slide,
         "Full run time: ~25–45 seconds  |  ~15,230 LLM tokens (automobile)  |  "
         "Sync path (CLI) or Async path (FastAPI)",
         Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
         size=Pt(13), color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — C++ TECHNICAL INDICATORS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Key Concept 7 — C++ Indicators via pybind11")
divider(slide, Inches(0.95))

text_box(slide,
         "RSI, MACD, and Bollinger Bands are implemented in C++ for correctness and speed.\n"
         "Python falls back gracefully if the .pyd extension isn't built.",
         Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.65), size=Pt(15), color=GRAY)

# Dispatch logic
code_block(slide,
           'try:\n'
           '    import stockindicators as _cpp_indicators\n'
           '    _USE_CPP = True     # C++ path\n'
           'except ImportError:\n'
           '    _USE_CPP = False    # pure-Python pandas fallback\n\n'
           'def compute_rsi(close, period=14):\n'
           '    if _USE_CPP:\n'
           '        return _cpp_indicators.compute_rsi(close.tolist(), period)\n'
           '    # pandas fallback:\n'
           '    delta = close.diff()\n'
           '    gain = delta.clip(lower=0)\n'
           '    loss = -delta.clip(upper=0)\n'
           '    avg_gain = gain.ewm(com=period-1, adjust=True).mean()\n'
           '    avg_loss = loss.ewm(com=period-1, adjust=True).mean()\n'
           '    rs = avg_gain / avg_loss\n'
           '    return 100 - (100 / (1 + rs))',
           Inches(0.5), Inches(1.8), Inches(6.5), Inches(4.6), font_size=Pt(11.5))

# Parity table
text_box(slide, "C++ ↔ Python Algorithm Parity", Inches(7.3), Inches(1.8),
         Inches(5.5), Inches(0.4), size=Pt(15), bold=True, color=CYAN)
two_col_table(slide,
    ["Indicator", "Python algo", "C++ impl", "Tolerance"],
    [
        ["RSI(14)", "ewm(com=13, adjust=True)", "ewm_adjusted(): recursive num/den", "±0.01"],
        ["MACD(12,26,9)", "ewm(span=n, adjust=False)", "ewm_unadjusted(): α·x+(1-α)·prev", "±0.01"],
        ["BB(20,2σ)", "rolling(20).std(ddof=1)", "Var=(Σx²−n·μ²)/(n−1)", "±0.01"],
    ],
    Inches(7.3), Inches(2.25),
    [Inches(1.2), Inches(1.7), Inches(1.9), Inches(0.8)],
    row_h=Inches(0.5), font_size=Pt(11))

text_box(slide, "Why C++?", Inches(7.3), Inches(4.45),
         Inches(5.5), Inches(0.4), size=Pt(15), bold=True, color=YELLOW)
text_box(slide,
         "• Exact EWM adjusted vs unadjusted matters —\n"
         "  pandas defaults differ for RSI vs MACD.\n"
         "  C++ encodes the correct algorithm explicitly.\n\n"
         "• pybind11 passes numpy arrays as C++ vectors\n"
         "  with zero Python overhead.\n\n"
         "• Build: powershell -File cpp/build_ext.ps1\n"
         "  Output: stockindicators.cp313-win_amd64.pyd",
         Inches(7.3), Inches(4.9), Inches(5.5), Inches(2.35), size=Pt(13), color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SAFETY RAILS & FALLBACKS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Safety Rails — No Single Failure Can Kill the Pipeline")
divider(slide, Inches(0.95))

# 3 rails
rails = [
    ("input_rail", "Before fan-out", "Validates ticker exists in yfinance.\nBad tickers flagged with rail_errors.\nAnalysis still continues.", CYAN),
    ("output_rail", "Inside run_agent", "Clamps every score to [0.0, 1.0].\nInjects default summary if empty.\nPrevents NaN from corrupting weighted sum.", YELLOW),
    ("conflict_rail", "Inside aggregate", "If any two agents diverge > 0.35,\ntriggers LLM re-resolution.\nEnsures final verdict reflects reasoning.", ORANGE),
]
for i, (name, when, desc, color) in enumerate(rails):
    bx = Inches(0.4 + i * 4.2)
    add_rect(slide, bx, Inches(1.1), Inches(3.9), Inches(2.3), BG_CARD, color, Pt(2))
    text_box(slide, name, bx + Inches(0.15), Inches(1.18), Inches(3.6), Inches(0.5),
             size=Pt(18), bold=True, color=color)
    text_box(slide, when, bx + Inches(0.15), Inches(1.68), Inches(3.6), Inches(0.38),
             size=Pt(13), italic=True, color=GRAY)
    text_box(slide, desc, bx + Inches(0.15), Inches(2.1), Inches(3.6), Inches(1.2),
             size=Pt(13), color=WHITE)

# Fallback chain table
text_box(slide, "Fallback Chain — what happens when each component fails",
         Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.4),
         size=Pt(15), bold=True, color=CYAN)
two_col_table(slide,
    ["Failure Point", "What Happens"],
    [
        ["Ticker resolution LLM fails", "Raw input used as ticker directly"],
        ["C++ .pyd absent (not built)", "_USE_CPP=False → pandas fallback, identical results"],
        ["RAG retrieval fails", "Falls back to Phase 2 live data (ContextBuilder)"],
        ["Live data fails (Serper/yfinance)", "Falls back to minimal stub string"],
        ["Agent LLM all retries fail", "AgentOutput(score=0.5) — neutral, never crashes pipeline"],
        ["Signal Aggregator JSON fail", "FinalReport(score=0.5, verdict=NEUTRAL)"],
        ["C# score_store proxy fails", "Warning logged → pipeline continues without persistence"],
    ],
    Inches(0.5), Inches(4.0),
    [Inches(4.2), Inches(7.8)],
    row_h=Inches(0.38), font_size=Pt(12))

text_box(slide,
         "Why 0.5 as fallback score?  It is the most conservative position — neither buy nor sell.\n"
         "A crashed agent producing a false BUY signal is far more dangerous than a neutral fallback.",
         Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.45),
         size=Pt(13), color=YELLOW)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ASYNC vs SYNC EXECUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Execution Paths — Sync vs Async")
divider(slide, Inches(0.95))

text_box(slide,
         "Two execution paths produce identical FinalReports. The path chosen depends on the caller.",
         Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5), size=Pt(15), color=GRAY)

# Left: sync
add_rect(slide, Inches(0.5), Inches(1.65), Inches(5.9), Inches(5.5), BG_CARD, GREEN, Pt(2))
text_box(slide, "SYNC PATH", Inches(0.65), Inches(1.75), Inches(5.6), Inches(0.45),
         size=Pt(18), bold=True, color=GREEN)
text_box(slide, "CLI · APScheduler · C# AnalyseJob",
         Inches(0.65), Inches(2.22), Inches(5.6), Inches(0.38),
         size=Pt(13), italic=True, color=GRAY)
code_block(slide,
           'ThreadPoolExecutor(max_workers=8)\n'
           '  ├── agent.run(query) × 8\n'
           '  └── as_completed(timeout=120s)\n\n'
           'LLM client: openai.OpenAI (sync)\n'
           'Retry back-off: time.sleep()\n\n'
           'Context I/O: blocking calls\n'
           '(yfinance, Serper) in threads',
           Inches(0.65), Inches(2.7), Inches(5.55), Inches(3.2), font_size=Pt(12))
text_box(slide, "Simple. No event loop. Works in any\nscript, scheduler, or test.", Inches(0.65), Inches(6.0), Inches(5.6), Inches(0.9),
         size=Pt(13), color=WHITE)

# Right: async
add_rect(slide, Inches(7.0), Inches(1.65), Inches(6.0), Inches(5.5), BG_CARD, CYAN, Pt(2))
text_box(slide, "ASYNC PATH", Inches(7.15), Inches(1.75), Inches(5.7), Inches(0.45),
         size=Pt(18), bold=True, color=CYAN)
text_box(slide, "FastAPI /analyse · WS /ws/stream",
         Inches(7.15), Inches(2.22), Inches(5.7), Inches(0.38),
         size=Pt(13), italic=True, color=GRAY)
code_block(slide,
           'asyncio.gather(*[\n'
           '    agent.run_async(query)\n'
           '    for agent in _SUB_AGENTS\n'
           '])\n\n'
           'each run_async():\n'
           '  context = await asyncio.to_thread(\n'
           '      _gather_context)  # I/O offloaded\n'
           '  raw = await AsyncOpenAI.chat\n'
           '      .completions.create()  # coroutine\n'
           'Retry back-off: asyncio.sleep()',
           Inches(7.15), Inches(2.7), Inches(5.7), Inches(3.2), font_size=Pt(12))
text_box(slide, "True coroutine-level concurrency.\nNo OS threads for LLM calls.\nLower memory under concurrent requests.", Inches(7.15), Inches(6.0), Inches(5.7), Inches(0.9),
         size=Pt(13), color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — KEYWORDS GLOSSARY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
heading(slide, "Keyword Glossary — Quick Reference")
divider(slide, Inches(0.95))

terms = [
    ("LLM Agent", "Prompt + context → LLM → structured JSON output. One agent = one viewpoint.", CYAN),
    ("LangGraph StateGraph", "Directed graph for LLM workflows. Nodes = functions, edges = flow.", YELLOW),
    ("Send()", "LangGraph fan-out: spawn N parallel branches simultaneously.", ORANGE),
    ("ContextBuilder", "Routes (sector, agent_name) to the right data fetcher (yfinance/Serper/Tavily).", GREEN),
    ("RAG", "Retrieval-Augmented Generation: query your own docs via vector similarity before LLM.", CYAN),
    ("ChromaDB", "Local vector database. Stores embeddings from PDFs/docs for RAG retrieval.", YELLOW),
    ("Signal Aggregator", "Weighted fusion + conflict detection + LLM resolution → final verdict.", ORANGE),
    ("Macro Cache", "TTL in-memory store. Background daemon pre-fetches sector news every 4h.", GREEN),
    ("input_rail / output_rail", "Guard nodes: validate ticker, clamp scores, inject defaults.", CYAN),
    ("pybind11", "C++ → Python bridge. Exposes RSI/MACD/BB as Python-callable functions.", YELLOW),
    ("Serper", "Google Search API. 2,500 free queries/month. Returns snippets, not full pages.", ORANGE),
    ("Tavily", "Full-page web extraction API. Used for policy PDFs (FAME, MNRE). 1,000/mo free.", GREEN),
    ("EWM (Exponential Weighted Mean)", "EWM adjusted (RSI) vs unadjusted (MACD) — different initial conditions.", CYAN),
    ("Conflict Resolution", "When two agents disagree by ≥0.30 on score, LLM reasons about which to trust.", YELLOW),
    ("Verdict", "STRONG BUY/BUY/NEUTRAL/SELL/STRONG SELL derived from 0–1 composite score.", ORANGE),
    ("FinalReport", "Pydantic model: score, verdict, sub-scores, thesis, risks, conflicts_resolved.", GREEN),
]

col1 = terms[:8]
col2 = terms[8:]
for i, (term, desc, color) in enumerate(col1):
    by = Inches(1.1 + i * 0.74)
    text_box(slide, f"  {term}", Inches(0.4), by, Inches(2.8), Inches(0.38),
             size=Pt(13), bold=True, color=color)
    text_box(slide, desc, Inches(3.2), by, Inches(3.0), Inches(0.38), size=Pt(12), color=WHITE)

for i, (term, desc, color) in enumerate(col2):
    by = Inches(1.1 + i * 0.74)
    text_box(slide, f"  {term}", Inches(6.7), by, Inches(3.2), Inches(0.38),
             size=Pt(13), bold=True, color=color)
    text_box(slide, desc, Inches(9.9), by, Inches(3.1), Inches(0.38), size=Pt(12), color=WHITE)

# vertical divider
add_rect(slide, Inches(6.45), Inches(1.0), Pt(2), Inches(6.2), GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
add_rect(slide, 0, 0, W, Inches(0.08), CYAN)

text_box(slide, "StockAgent in One Sentence",
         Inches(0.6), Inches(0.8), Inches(12), Inches(0.55),
         size=Pt(22), bold=True, color=GRAY)

text_box(slide,
         '"28 specialist LLM agents, each given real-time data for their domain,\n'
         ' run in parallel via LangGraph, and their scores are fused with\n'
         ' conflict-aware weighted aggregation into a single investment verdict."',
         Inches(0.6), Inches(1.4), Inches(12), Inches(1.6),
         size=Pt(22), bold=True, color=WHITE)

divider(slide, Inches(3.1), CYAN)

text_box(slide, "What makes it interesting:", Inches(0.6), Inches(3.3),
         Inches(12), Inches(0.4), size=Pt(18), bold=True, color=CYAN)

highlights = [
    "🔀  Conflict resolution uses LLM reasoning, not arithmetic — captures which signal matters MORE right now",
    "💾  Macro Cache: background daemon pre-fetches sector news so individual analyses are cache-HIT (saves ~1,000 Serper calls/month)",
    "🏗️   4-language system: C++ computes indicators, Python runs agents, TypeScript serves the UI, C# schedules and persists",
    "🛡️   Three safety rails + full fallback chain — no single failure can crash the pipeline or produce a false signal",
    "📐  ContextBuilder two-level routing: sector-specific builders fall through to shared yfinance builders (pattern_analysis works for all sectors)",
]
for i, h in enumerate(highlights):
    text_box(slide, h, Inches(0.6), Inches(3.8 + i * 0.62), Inches(12.2), Inches(0.55),
             size=Pt(14), color=WHITE)

text_box(slide, "Questions?", Inches(0.6), Inches(7.0), Inches(12), Inches(0.4),
         size=Pt(28), bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), CYAN)


# ─── Save ─────────────────────────────────────────────────────────────────────
out_path = Path("outputs/StockAgent_Presentation.pptx")
out_path.parent.mkdir(exist_ok=True)
prs.save(str(out_path))
print(f"DONE  Saved: {out_path}  ({out_path.stat().st_size // 1024} KB,  {len(prs.slides)} slides)")
